"""Generate PowerPoint presentation for Editor Recommender Memory System."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0x9D, 0xC3, 0xE6)
ACCENT_GREEN = RGBColor(0x00, 0xB0, 0x50)
ACCENT_ORANGE = RGBColor(0xED, 0x7D, 0x31)
ACCENT_RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x59, 0x56, 0x59)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=BLACK, bold_first=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
        # bullet
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        # remove existing bullets
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BLUE)

add_text_box(slide, 1.5, 1.5, 10, 1.2,
             "Editor Recommender — Memory System",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.0, 10, 0.8,
             "From Stateless to Learning: Session & Long-Term Memory POC",
             font_size=22, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 4.5, 10, 0.5,
             "February 2026  |  v1.3.0  |  Dev Environment",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.5, 10, 0.5,
             "GTS AI Team",
             font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "The Problem: A Stateless System",
             font_size=28, bold=True, color=DARK_BLUE)

# Before box
add_rounded_rect(slide, 0.8, 1.5, 5.5, 4.8, LIGHT_GRAY, "", font_size=10)
add_text_box(slide, 1.0, 1.6, 5.0, 0.5,
             "Before (Stateless)", font_size=20, bold=True, color=ACCENT_RED)
add_bullet_list(slide, 1.0, 2.3, 5.0, 3.5, [
    "Every workflow started from scratch — no memory of past decisions",
    "If a pod crashed mid-workflow, all progress was lost",
    "LLM made the same analysis repeatedly for similar manuscripts",
    "No audit trail — couldn't trace why an editor was picked",
    "No way to learn from patterns in past assignments",
], font_size=14, color=GRAY)

# After box
add_rounded_rect(slide, 7.0, 1.5, 5.5, 4.8, LIGHT_GRAY, "", font_size=10)
add_text_box(slide, 7.2, 1.6, 5.0, 0.5,
             "After (With Memory)", font_size=20, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, 7.2, 2.3, 5.0, 3.5, [
    "Every decision is saved to PostgreSQL — survives pod restarts",
    "Crash recovery: workflow resumes from last checkpoint",
    "LLM sees past assignments as context for new decisions",
    "Full audit trail of every step in every workflow",
    "System learns and improves with each manuscript processed",
], font_size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Memory Architecture: Two Tiers",
             font_size=28, bold=True, color=DARK_BLUE)

# Tier 2 Box
add_rounded_rect(slide, 0.8, 1.5, 5.5, 5.0, RGBColor(0xE8, 0xF0, 0xFE), "")
add_text_box(slide, 1.0, 1.6, 5.0, 0.5,
             "Tier 2 — Session Memory", font_size=20, bold=True, color=MEDIUM_BLUE)
add_text_box(slide, 1.0, 2.2, 5.0, 0.4,
             "LangGraph AsyncPostgresSaver", font_size=13, bold=False, color=GRAY)
add_bullet_list(slide, 1.0, 2.7, 5.0, 3.5, [
    "Saves workflow state after EACH graph node executes",
    "If pod crashes → resumes from last checkpoint",
    "Thread-based: each manuscript gets unique thread_id",
    "Full audit trail of every step",
    "Tables: checkpoints, checkpoint_blobs, checkpoint_writes",
], font_size=13, color=GRAY)

# Tier 3 Box
add_rounded_rect(slide, 7.0, 1.5, 5.5, 5.0, RGBColor(0xE2, 0xF0, 0xD9), "")
add_text_box(slide, 7.2, 1.6, 5.0, 0.5,
             "Tier 3 — Long-Term Memory", font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.2, 2.2, 5.0, 0.4,
             "LangGraph AsyncPostgresStore", font_size=13, bold=False, color=GRAY)
add_bullet_list(slide, 7.2, 2.7, 5.0, 3.5, [
    "Saves completed assignment decisions permanently",
    "Organized by journal: namespace = (\"assignments\", journal_id)",
    "Each entry stores: editor, reasoning, runner-up, topics",
    "Before each new decision, system queries past assignments",
    "Past assignments injected into LLM prompt as context",
], font_size=13, color=GRAY)

# Database label
add_rounded_rect(slide, 3.5, 6.6, 6.3, 0.6, DARK_BLUE,
                 "PostgreSQL 17.4  |  mspubs-dev RDS  |  mspubs schema  |  6 tables",
                 font_size=13, font_color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Database Tables & Schema
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Database: Tables & Data Fields",
             font_size=28, bold=True, color=DARK_BLUE)

# Session Memory tables
add_text_box(slide, 0.8, 1.4, 5.5, 0.4,
             "Session Memory Tables (Tier 2)", font_size=18, bold=True, color=MEDIUM_BLUE)

table_data_session = [
    ["Table", "Key Columns", "Rows"],
    ["checkpoints", "thread_id, checkpoint_id, type, checkpoint (jsonb), metadata (jsonb)", "42"],
    ["checkpoint_blobs", "thread_id, channel, version, blob (bytea)", "22"],
    ["checkpoint_writes", "thread_id, checkpoint_id, task_id, channel, blob (bytea)", "76"],
    ["checkpoint_migrations", "v (integer)", "10"],
]

tbl = slide.shapes.add_table(len(table_data_session), 3,
                              Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.2)).table
for i, row in enumerate(table_data_session):
    for j, cell_text in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = MEDIUM_BLUE

tbl.columns[0].width = Inches(1.5)
tbl.columns[1].width = Inches(3.2)
tbl.columns[2].width = Inches(0.8)

# Long-term Memory tables
add_text_box(slide, 7.0, 1.4, 5.5, 0.4,
             "Long-Term Memory Tables (Tier 3)", font_size=18, bold=True, color=ACCENT_GREEN)

table_data_store = [
    ["Table", "Key Columns", "Rows"],
    ["store", "prefix, key, value (jsonb), created_at, updated_at", "5"],
    ["store_migrations", "v (integer)", "4"],
]

tbl2 = slide.shapes.add_table(len(table_data_store), 3,
                               Inches(7.0), Inches(1.9), Inches(5.5), Inches(1.0)).table
for i, row in enumerate(table_data_store):
    for j, cell_text in enumerate(row):
        cell = tbl2.cell(i, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_GREEN

tbl2.columns[0].width = Inches(1.5)
tbl2.columns[1].width = Inches(3.2)
tbl2.columns[2].width = Inches(0.8)

# Store value schema
add_text_box(slide, 7.0, 3.2, 5.5, 0.4,
             "Store Value Schema (what we save per assignment):",
             font_size=14, bold=True, color=DARK_BLUE)

add_bullet_list(slide, 7.0, 3.7, 5.5, 3.0, [
    "editor_id — ORCID of selected editor",
    "editor_person_id — PersonId of selected editor",
    "reasoning — LLM's full reasoning for the decision",
    "runner_up — Second-choice editor",
    "filtered_out_editors — Editors removed due to COI",
    "journal_id — Which journal (jm, ja, oc, etc.)",
    "manuscript_number — Unique manuscript identifier",
    "topics — Extracted research topics from reasoning",
    "timestamp — When the decision was made",
], font_size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Learning Loop Flow
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "How the Learning Loop Works",
             font_size=28, bold=True, color=DARK_BLUE)

# Step boxes
steps = [
    ("1. Manuscript Arrives", "POST /execute_workflow\nwith manuscript_number & journal_id", MEDIUM_BLUE, 0.5),
    ("2. Fetch Data", "EE API returns manuscript info\n+ available editors from Oracle DB", MEDIUM_BLUE, 2.0),
    ("3. Search Memory", "Query store table:\n\"Any past decisions for this journal?\"", ACCENT_GREEN, 3.5),
    ("4. Inject Into Prompt", "Past assignments formatted and added\nto LLM prompt as reference context", ACCENT_GREEN, 5.0),
]

for title, desc, color, top in steps:
    add_rounded_rect(slide, 0.8, top, 3.5, 1.2, color, title, font_size=14, font_color=WHITE)
    add_text_box(slide, 4.5, top + 0.1, 3.5, 1.0, desc, font_size=12, color=GRAY)

steps2 = [
    ("5. LLM Decides", "Claude reads: editors + manuscript +\njournal rules + PAST DECISIONS", MEDIUM_BLUE, 0.5),
    ("6. Save to Memory", "Decision saved to store BEFORE\ncalling assignment API", ACCENT_GREEN, 2.0),
    ("7. Call Assignment API", "Submit editor assignment\nto downstream system", MEDIUM_BLUE, 3.5),
    ("8. Next Run Benefits", "Future manuscripts see this decision\nas reference → better consistency", ACCENT_ORANGE, 5.0),
]

for title, desc, color, top in steps2:
    add_rounded_rect(slide, 8.0, top, 3.5, 1.2, color, title, font_size=14, font_color=WHITE)
    add_text_box(slide, 8.0, top + 1.3, 3.5, 0.3, desc, font_size=12, color=GRAY)

# Key insight
add_rounded_rect(slide, 0.8, 6.5, 11.7, 0.7, RGBColor(0xFF, 0xF2, 0xCC),
                 "Key: Save happens BEFORE the assignment API call — even if downstream fails, the LLM's decision is captured",
                 font_size=13, font_color=DARK_BLUE, bold=False)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Real Example: Feb 24 (First Decision)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Real Example: Feb 24 — First Decision (No Memory Yet)",
             font_size=28, bold=True, color=DARK_BLUE)

# Timeline marker
add_rounded_rect(slide, 0.8, 1.4, 2.0, 0.5, ACCENT_ORANGE,
                 "Feb 24, 2026", font_size=14)

add_text_box(slide, 3.2, 1.4, 8.0, 0.5,
             "Manuscript: jm-2021-01708p  |  Journal: jm  |  Store: EMPTY (0 past decisions)",
             font_size=14, bold=True, color=GRAY)

# What LLM received
add_rounded_rect(slide, 0.8, 2.2, 5.5, 4.5, RGBColor(0xE8, 0xF0, 0xFE), "")
add_text_box(slide, 1.0, 2.3, 5.0, 0.4,
             "What LLM Received in Prompt:", font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_list(slide, 1.0, 2.8, 5.0, 3.5, [
    "Manuscript Type: Article, Peer-Reviewed",
    "Author Institutions: None provided",
    "Editor 1: 130958 (Editor-in-Chief, rank 1, Vanderbilt)",
    "Editor 2: 971391 (Executive Editor, rank 2, Oxford)",
    "Journal Rules: PATH C = pick lowest rank",
    "Past Assignments: NONE ← empty memory",
], font_size=13, color=GRAY)

# What LLM decided
add_rounded_rect(slide, 7.0, 2.2, 5.5, 4.5, RGBColor(0xE2, 0xF0, 0xD9), "")
add_text_box(slide, 7.2, 2.3, 5.0, 0.4,
             "LLM Decision:", font_size=16, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, 7.2, 2.8, 5.0, 1.5, [
    "Selected Editor: 130958 (Editor-in-Chief)",
    "Reasoning: PATH C — rank 1 is lowest available",
    "Runner-up: 971391 (Executive Editor, rank 2)",
    "COI: None detected",
], font_size=13, color=GRAY)

add_text_box(slide, 7.2, 4.5, 5.0, 0.4,
             "What Got Saved to Memory:", font_size=16, bold=True, color=DARK_BLUE)
add_bullet_list(slide, 7.2, 5.0, 5.0, 1.5, [
    "Namespace: (\"assignments\", \"jm\")",
    "Key: jm-2021-01708p",
    "Value: editor 130958, reasoning, runner-up, topics",
], font_size=13, color=GRAY)

# Pod log evidence
add_rounded_rect(slide, 0.8, 6.8, 11.7, 0.5, RGBColor(0x1B, 0x1B, 0x1B),
                 "Pod Log: \"Saved assignment to long-term memory: jm-2021-01708p → editor 130958\"",
                 font_size=12, font_color=ACCENT_GREEN, bold=False)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Real Example: Feb 26 (Memory Kicks In)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Real Example: Feb 26 — Second Decision (Memory Active!)",
             font_size=28, bold=True, color=DARK_BLUE)

# Timeline marker
add_rounded_rect(slide, 0.8, 1.4, 2.0, 0.5, ACCENT_GREEN,
                 "Feb 26, 2026", font_size=14)

add_text_box(slide, 3.2, 1.4, 9.0, 0.5,
             "Manuscript: jm-2021-018697  |  Journal: jm  |  Store: 1 past decision for jm ✓",
             font_size=14, bold=True, color=GRAY)

# What happened step by step
add_rounded_rect(slide, 0.8, 2.2, 3.8, 4.5, RGBColor(0xE8, 0xF0, 0xFE), "")
add_text_box(slide, 1.0, 2.3, 3.5, 0.4,
             "Step 1: Memory Search", font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_list(slide, 1.0, 2.8, 3.5, 3.5, [
    "System queries store table",
    "Filter: namespace = (assignments, jm)",
    "Found 1 result: jm-2021-01708p",
    "The Feb 24 decision!",
], font_size=12, color=GRAY)

add_rounded_rect(slide, 4.8, 2.2, 3.8, 4.5, RGBColor(0xE2, 0xF0, 0xD9), "")
add_text_box(slide, 5.0, 2.3, 3.5, 0.4,
             "Step 2: Prompt Injection", font_size=16, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, 5.0, 2.8, 3.5, 3.5, [
    "Past assignment formatted as text:",
    "\"Manuscript jm-2021-01708p\"",
    "\"Assigned Editor: 130958\"",
    "\"Reasoning: PATH C, rank 1\"",
    "Injected into {past_assignments}",
    "placeholder in LLM prompt",
], font_size=12, color=GRAY)

add_rounded_rect(slide, 8.8, 2.2, 3.8, 4.5, RGBColor(0xFC, 0xE4, 0xD6), "")
add_text_box(slide, 9.0, 2.3, 3.5, 0.4,
             "Step 3: LLM Decides", font_size=16, bold=True, color=ACCENT_ORANGE)
add_bullet_list(slide, 9.0, 2.8, 3.5, 3.5, [
    "LLM sees current manuscript",
    "LLM sees available editors",
    "LLM sees journal rules",
    "LLM sees past decision ← NEW!",
    "Decides: editor 130958",
    "Consistent with prior pattern",
], font_size=12, color=GRAY)

# Pod log evidence
add_rounded_rect(slide, 0.8, 6.8, 11.7, 0.5, RGBColor(0x1B, 0x1B, 0x1B),
                 "Pod Log: \"Long-term memory search returned 1 results\" → \"injecting 1 similar past assignments into prompt\"",
                 font_size=12, font_color=ACCENT_GREEN, bold=False)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Data Growth Evidence
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Data Growth: Before & After",
             font_size=28, bold=True, color=DARK_BLUE)

# Before/After table
table_data = [
    ["Table", "Feb 24", "Feb 26", "Change"],
    ["checkpoints", "42", "52+", "↑ New workflow checkpoints"],
    ["checkpoint_writes", "76", "88+", "↑ New state writes"],
    ["store", "4", "5", "↑ 1 new real assignment"],
    ["Threads tracked", "5", "6", "↑ jm-2021-018697 added"],
]

tbl3 = slide.shapes.add_table(len(table_data), 4,
                               Inches(0.8), Inches(1.5), Inches(7.0), Inches(2.5)).table
for i, row in enumerate(table_data):
    for j, cell_text in enumerate(row):
        cell = tbl3.cell(i, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif j == 3:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True
            else:
                p.font.color.rgb = GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

tbl3.columns[0].width = Inches(1.5)
tbl3.columns[1].width = Inches(1.0)
tbl3.columns[2].width = Inches(1.0)
tbl3.columns[3].width = Inches(3.5)

# Store contents
add_text_box(slide, 0.8, 4.3, 11, 0.4,
             "Current Store Contents (Long-Term Memory):", font_size=18, bold=True, color=DARK_BLUE)

store_data = [
    ["Namespace / Key", "Editor", "Updated", "Source"],
    ["assignments.jm / jm-2021-018697", "130958", "Feb 26, 13:00", "NEW — 2nd workflow run"],
    ["assignments.jm / jm-2021-01708p", "130958", "Feb 24, 14:43", "1st real workflow"],
    ["editor_assignments.ja / ja-2025-test003", "ED003", "Feb 23, 08:16", "Test data"],
    ["editor_assignments.ja / ja-2025-test002", "ED002", "Feb 23, 08:16", "Test data"],
    ["editor_assignments.ja / ja-2025-test001", "ED001", "Feb 23, 08:16", "Test data"],
]

tbl4 = slide.shapes.add_table(len(store_data), 4,
                               Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.5)).table
for i, row in enumerate(store_data):
    for j, cell_text in enumerate(row):
        cell = tbl4.cell(i, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif i == 1:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True
            else:
                p.font.color.rgb = GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

tbl4.columns[0].width = Inches(4.5)
tbl4.columns[1].width = Inches(1.5)
tbl4.columns[2].width = Inches(2.5)
tbl4.columns[3].width = Inches(3.2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Files Changed
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Implementation: Files & Tests",
             font_size=28, bold=True, color=DARK_BLUE)

files_data = [
    ["File", "Lines", "What It Does"],
    ["memory.py (NEW)", "367", "Core memory module — create_checkpointer(), create_store(), save_assignment_to_memory(), search_similar_assignments(), format_past_assignments_for_prompt()"],
    ["ee_graph_anthropic.py", "472", "Integrated memory read (search past assignments before LLM call) and memory write (save decision before assignment API)"],
    ["app.py", "220", "FastAPI lifespan initializes memory when POSTGRES_URI is set; graceful degradation without it"],
    ["test_memory.py (NEW)", "336", "6 unit tests against real Postgres — checkpointer, store, save, search, format"],
    ["test_memory_integration.py (NEW)", "490", "4 integration tests — full workflow with mocked APIs, verifying session + long-term memory"],
]

tbl5 = slide.shapes.add_table(len(files_data), 3,
                               Inches(0.8), Inches(1.3), Inches(11.7), Inches(3.5)).table
for i, row in enumerate(files_data):
    for j, cell_text in enumerate(row):
        cell = tbl5.cell(i, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

tbl5.columns[0].width = Inches(2.5)
tbl5.columns[1].width = Inches(0.7)
tbl5.columns[2].width = Inches(8.5)

# Tech stack
add_text_box(slide, 0.8, 5.2, 11, 0.4,
             "Technology Stack:", font_size=18, bold=True, color=DARK_BLUE)
add_bullet_list(slide, 0.8, 5.7, 11, 1.5, [
    "langgraph-checkpoint-postgres v3.0.4 — AsyncPostgresSaver for session memory",
    "langgraph-store-postgres — AsyncPostgresStore for long-term key-value memory (pgvector-ready)",
    "psycopg[binary] v3.2 + AsyncConnectionPool — async Postgres connections",
    "PostgreSQL 17.4 on AWS RDS (mspubs-dev) — existing infrastructure, no new database needed",
], font_size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — What's Next
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "What's Next: Roadmap",
             font_size=28, bold=True, color=DARK_BLUE)

# Done
add_rounded_rect(slide, 0.8, 1.5, 3.8, 5.0, RGBColor(0xE2, 0xF0, 0xD9), "")
add_text_box(slide, 1.0, 1.6, 3.5, 0.4,
             "✓ Completed (POC)", font_size=18, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, 1.0, 2.2, 3.5, 4.0, [
    "Session memory (checkpoints)",
    "Long-term memory (key-value store)",
    "Memory save before API call",
    "Prompt injection of past decisions",
    "10 tests (6 unit + 4 integration)",
    "Deployed v1.3.0 to dev",
    "Real data captured & verified",
    "Learning loop proven end-to-end",
], font_size=12, color=GRAY)

# In Progress
add_rounded_rect(slide, 4.8, 1.5, 3.8, 5.0, RGBColor(0xFF, 0xF2, 0xCC), "")
add_text_box(slide, 5.0, 1.6, 3.5, 0.4,
             "→ Next Phase", font_size=18, bold=True, color=ACCENT_ORANGE)
add_bullet_list(slide, 5.0, 2.2, 3.5, 4.0, [
    "pgvector for semantic search",
    "Bedrock Titan embeddings",
    "Topic-based manuscript matching",
    "Scale test with more manuscripts",
    "Journals with 10+ editors",
    "Human feedback integration",
    "Memory cleanup / TTL policies",
], font_size=12, color=GRAY)

# Future
add_rounded_rect(slide, 8.8, 1.5, 3.8, 5.0, RGBColor(0xE8, 0xF0, 0xFE), "")
add_text_box(slide, 9.0, 1.6, 3.5, 0.4,
             "◇ Future Vision", font_size=18, bold=True, color=MEDIUM_BLUE)
add_bullet_list(slide, 9.0, 2.2, 3.5, 4.0, [
    "Learn from editorial outcomes",
    "Track acceptance/rejection rates",
    "Editor workload balancing",
    "Cross-journal pattern sharing",
    "Confidence scoring over time",
    "Human-in-the-loop approval",
    "A/B testing memory vs no-memory",
], font_size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Summary
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)

add_text_box(slide, 1.5, 1.0, 10, 0.8,
             "Summary",
             font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bullet_list(slide, 2.0, 2.2, 9.0, 4.5, [
    "Built a 2-tier memory system that makes the Editor Recommender learn from past decisions",
    "Session Memory: Every workflow step is checkpointed — crash recovery + full audit trail",
    "Long-Term Memory: Every assignment decision is saved and used as context for future runs",
    "Proven with real data: Feb 24 decision was found and injected into Feb 26 prompt",
    "The system now gets smarter with every manuscript processed",
    "Next: Semantic search (pgvector) for topic-aware matching across manuscripts",
], font_size=18, color=WHITE, spacing=Pt(14))

add_text_box(slide, 1.5, 6.5, 10, 0.5,
             "v1.3.0  |  10 Tests Passing  |  5 Store Entries  |  6 Threads Tracked  |  Dev Deployed",
             font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
output_path = r"c:\Users\szs17\OneDrive - ACS\Documents\GitHub\at-ai-editor-recommender\docs\Editor_Recommender_Memory_System.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
