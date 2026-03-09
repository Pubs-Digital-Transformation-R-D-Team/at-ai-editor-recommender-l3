"""
Build: Strands Agents Migration PPT
Run:   python docs/build_strands_ppt.py
Output: docs/Strands_Migration_Editor_Recommender.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
AWS_ORANGE   = RGBColor(0xFF, 0x99, 0x00)   # AWS orange
DARK_NAVY    = RGBColor(0x1A, 0x23, 0x3A)   # slide backgrounds
MID_BLUE     = RGBColor(0x23, 0x4F, 0x8C)   # section headers
LIGHT_BLUE   = RGBColor(0x4A, 0x9E, 0xD6)   # accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF2, 0xF4, 0xF8)
GREEN        = RGBColor(0x27, 0xAE, 0x60)
RED          = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW       = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color=DARK_NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def accent_bar(slide, color=AWS_ORANGE, t=0.0, h=0.07):
    box(slide, 0, t, 13.33, h, fill_color=color)

def slide_header(slide, title, subtitle=None):
    accent_bar(slide, AWS_ORANGE, 0.0, 0.07)
    txt(slide, title, 0.4, 0.15, 12.5, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.75, 12.5, 0.45, size=16, color=LIGHT_BLUE)
    accent_bar(slide, MID_BLUE, 7.38, 0.07)

def bullet_block(slide, items, l, t, w, h, title=None, title_color=AWS_ORANGE,
                 item_size=15, title_size=17):
    y = t
    if title:
        txt(slide, title, l, y, w, 0.35, size=title_size, bold=True,
            color=title_color)
        y += 0.38
    for icon, item in items:
        txt(slide, f"{icon}  {item}", l, y, w, 0.32, size=item_size, color=WHITE)
        y += 0.33

def card(slide, l, t, w, h, heading, lines, head_color=MID_BLUE,
         bg_color=RGBColor(0x1E, 0x2D, 0x4A), line_size=13):
    box(slide, l, t, w, h, fill_color=bg_color,
        line_color=head_color, line_width=Pt(1.5))
    txt(slide, heading, l+0.1, t+0.08, w-0.2, 0.35,
        size=14, bold=True, color=head_color)
    y = t + 0.42
    for line in lines:
        txt(slide, line, l+0.15, y, w-0.25, 0.28, size=line_size, color=WHITE)
        y += 0.28


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
accent_bar(s, AWS_ORANGE, 0.0, 0.12)
accent_bar(s, MID_BLUE, 7.38, 0.12)

box(s, 0.5, 1.1, 12.33, 4.8, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
    line_color=AWS_ORANGE, line_width=Pt(2))

txt(s, "AWS Strands Agents SDK", 0.9, 1.4, 11.5, 1.0,
    size=40, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
txt(s, "Editor Recommender — Memory & Workflow Migration", 0.9, 2.35, 11.5, 0.7,
    size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "LangGraph → Strands  |  v1.3.0 → v1.4.0  |  March 2026", 0.9, 3.05,
    11.5, 0.5, size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txt(s, "Branch: feature/strands-memory-implementation", 0.9, 4.0,
    11.5, 0.4, size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
txt(s, "AWS Account: 412381768680  |  Region: us-east-1  |  EKS Namespace: er",
    0.9, 4.4, 11.5, 0.4, size=13,
    color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

txt(s, "GTS AI / Pubs Engineering", 0.9, 5.2, 11.5, 0.4,
    size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Agenda")

items = [
    ("1", "Why Migrate? — LangGraph limitations & Strands advantages"),
    ("2", "Architecture: LangGraph vs Strands side-by-side"),
    ("3", "Memory System: 3-Tier design (unchanged Tier 3 + new Tier 2)"),
    ("4", "Code Changes — what was added, modified, deleted"),
    ("5", "Live Data: Learning loop proof (Feb 24 → Feb 26 → Mar 2)"),
    ("6", "Workaround: In-memory session fallback (no S3 IRSA needed)"),
    ("7", "Dev Deployment Checklist"),
    ("8", "Next Steps"),
]
y = 1.3
for num, item in items:
    box(s, 0.5, y, 0.45, 0.38, fill_color=AWS_ORANGE)
    txt(s, num, 0.5, y, 0.45, 0.38, size=16, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, item, 1.05, y+0.03, 11.5, 0.38, size=15, color=WHITE)
    y += 0.50


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why Migrate?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Why Migrate from LangGraph to Strands?")

# LangGraph problems
card(s, 0.4, 1.2, 5.8, 5.6, "❌  LangGraph Pain Points",
     [
         "• Explicit graph: nodes + edges hard-coded",
         "• 3 workflow files (1,094 lines combined)",
         "• AsyncPostgresSaver → 6 checkpoint tables in RDS",
         "• Anthropic SDK dependency (separate from AWS)",
         "• langchain-aws, langchain-mcp-adapters required",
         "• Graph state machine verbose & hard to extend",
         "• langgraph-cli, langgraph-checkpoint-postgres",
         "• Tight coupling: tool calls in fixed sequence",
         "• Tier 2 (Postgres) adds latency + RDS storage cost",
     ],
     head_color=RED, line_size=13)

# Strands benefits
card(s, 6.6, 1.2, 6.3, 5.6, "✅  Strands Advantages",
     [
         "• Model-driven: LLM decides which tools to call",
         "• Single file: ee_agent_strands.py (~500 lines)",
         "• S3SessionManager → 1 JSON object per run",
         "• Native AWS (boto3/IRSA — no new K8s changes)",
         "• Only dependency: strands-agents>=1.29.0",
         "• @tool decorator: simple, readable, testable",
         "• Removes langgraph, langchain-aws, anthropic",
         "• Flexible: LLM skips tools it doesn't need",
         "• Tier 2 (S3) cheaper, simpler, serverless",
     ],
     head_color=GREEN, line_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Comparison
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Architecture: LangGraph vs Strands")

# LangGraph diagram
box(s, 0.3, 1.15, 5.9, 5.9, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
    line_color=RED, line_width=Pt(1.5))
txt(s, "BEFORE — LangGraph", 0.5, 1.2, 5.5, 0.4, size=15, bold=True, color=RED)

lg_steps = [
    ("app.py", "FastAPI  →  checkpointer + store init", MID_BLUE),
    ("ee_graph_anthropic.py", "Graph: START node", MID_BLUE),
    ("", "→  fetch_manuscript_data node", LIGHT_BLUE),
    ("", "→  search_long_term_memory node", LIGHT_BLUE),
    ("", "→  call_llm node (Anthropic Bedrock)", LIGHT_BLUE),
    ("", "→  parse_result node", LIGHT_BLUE),
    ("", "→  save_to_memory node", LIGHT_BLUE),
    ("", "→  call_assign_api node", LIGHT_BLUE),
    ("", "→  END node", LIGHT_BLUE),
    ("memory.py", "AsyncPostgresSaver (Tier 2)", YELLOW),
    ("", "AsyncPostgresStore (Tier 3)", YELLOW),
]
y = 1.7
for fname, step, color in lg_steps:
    prefix = f"[{fname}] " if fname else "   "
    txt(s, f"{prefix}{step}", 0.5, y, 5.6, 0.28, size=11, color=color)
    y += 0.29

# Arrow
txt(s, "→", 6.3, 3.6, 0.7, 0.5, size=30, bold=True,
    color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Strands diagram
box(s, 7.1, 1.15, 5.9, 5.9, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
    line_color=GREEN, line_width=Pt(1.5))
txt(s, "AFTER — Strands", 7.3, 1.2, 5.5, 0.4, size=15, bold=True, color=GREEN)

st_steps = [
    ("app.py", "FastAPI  →  store init only (no checkpointer)", MID_BLUE),
    ("ee_agent_strands.py", "EditorAssignmentAgent", MID_BLUE),
    ("", "  BedrockModel (Nova Premier)", LIGHT_BLUE),
    ("", "  @tool: fetch_manuscript_data()", LIGHT_BLUE),
    ("", "  @tool: search_past_assignments(query)", LIGHT_BLUE),
    ("", "  S3SessionManager (Tier 2)", YELLOW),
    ("", "  Agent.invoke_async(task)", GREEN),
    ("", "    LLM calls tools in any order it decides", LIGHT_BLUE),
    ("", "    LLM returns JSON recommendation", LIGHT_BLUE),
    ("", "  save_to_long_term_memory() → Postgres", YELLOW),
    ("", "  call_assign_api()", MID_BLUE),
    ("memory.py", "AsyncPostgresStore (Tier 3 — unchanged)", YELLOW),
]
y = 1.7
for fname, step, color in st_steps:
    prefix = f"[{fname}] " if fname else "   "
    txt(s, f"{prefix}{step}", 7.3, y, 5.6, 0.28, size=11, color=color)
    y += 0.29


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Memory Architecture (3-tier)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "3-Tier Memory Architecture", "What changed and what stayed the same")

# Tier 1
card(s, 0.4, 1.2, 3.9, 2.7, "Tier 1 — Working Memory",
     [
         "In-process Python state",
         "Lives only for one request",
         "Variables: manuscript_data,",
         "  editors_list, llm_output",
         "No persistence — same as before",
     ], head_color=LIGHT_BLUE, line_size=13)
txt(s, "UNCHANGED", 0.85, 3.6, 3.0, 0.35, size=13, bold=True,
    color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Tier 2
card(s, 4.7, 1.2, 4.0, 2.7, "Tier 2 — Session Memory",
     [
         "BEFORE: AsyncPostgresSaver",
         "  → 6 checkpoint tables in RDS",
         "  → Postgres dependency",
         "AFTER: S3SessionManager",
         "  → 1 JSON per manuscript run",
         "  → bucket: acs-us-east-1-gtsai-",
         "      dev-s3-er-sessions",
     ], head_color=AWS_ORANGE, line_size=12)
txt(s, "REPLACED", 5.2, 3.6, 3.0, 0.35, size=13, bold=True,
    color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Tier 3
card(s, 9.1, 1.2, 3.9, 2.7, "Tier 3 — Long-term Memory",
     [
         "AsyncPostgresStore",
         "DB: mspubs-dev RDS (mspubs schema)",
         "Tables: store + writes",
         "5 rows (Feb 24 – Mar 2, 2026)",
         "Semantic search via pgvector",
         "Learning loop: PROVEN",
     ], head_color=GREEN, line_size=12)
txt(s, "UNCHANGED ✅", 9.5, 3.6, 3.0, 0.35, size=13, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)

# Flow arrows
for x in [4.3, 8.7]:
    txt(s, "→", x, 2.3, 0.5, 0.4, size=20, bold=True,
        color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Bottom data box
box(s, 0.4, 4.1, 12.53, 2.9, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
    line_color=LIGHT_BLUE, line_width=Pt(1))
txt(s, "Live Tier 3 Data (Mar 2, 2026 snapshot)", 0.6, 4.15, 12.0, 0.35,
    size=14, bold=True, color=LIGHT_BLUE)

data_cols = [
    ("Postgres Table", "mspubs.checkpoint", "mspubs.checkpoint_blobs",
     "mspubs.checkpoint_migrations", "mspubs.checkpoint_writes",
     "mspubs.store", "mspubs.store_migrations"),
    ("Row Count", "48", "96", "1", "0", "5 (learning data)", "1"),
    ("Purpose", "LangGraph checkpoints (retired)", "Blob storage (retired)",
     "Schema version (retired)", "Write ops (retired)",
     "Long-term assignments ← ACTIVE", "Store version"),
]
col_widths = [3.8, 2.2, 5.8]
col_x = [0.6, 4.6, 7.0]
for ci, (col, cx, cw) in enumerate(zip(data_cols, col_x, col_widths)):
    y = 4.55
    for ri, cell in enumerate(col):
        cell_color = WHITE if ri == 0 else (
            GREEN if "ACTIVE" in cell else
            RGBColor(0xAA, 0xAA, 0xAA) if "retired" in cell else WHITE)
        txt(s, cell, cx, y, cw, 0.28, size=11,
            bold=(ri == 0), color=cell_color)
        y += 0.28


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Code Changes
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Code Changes Summary")

# DELETED
card(s, 0.4, 1.2, 3.8, 4.8, "🗑  DELETED",
     [
         "ee_graph_anthropic.py",
         "  472 lines — LangGraph workflow",
         "  nodes/edges/state machine",
         "",
         "ee_graph_claude.py",
         "  ~200 lines — unused variant",
         "",
         "ee_graph.py",
         "  ~150 lines — original file",
         "",
         "ee_workflow_langgraph_otel.py",
         "  script entry (stale reference)",
     ], head_color=RED, line_size=12)

# MODIFIED
card(s, 4.6, 1.2, 4.3, 4.8, "✏️  MODIFIED",
     [
         "app.py",
         "  • Strands-only (no toggle)",
         "  • lifespan: store only",
         "  • ConfigDict (Pydantic V2)",
         "",
         "memory.py",
         "  • Removed create_checkpointer()",
         "  • Tier 3 unchanged",
         "",
         "client.py",
         "  • Import from ee_agent_strands",
         "",
         "pyproject.toml",
         "  • Removed 4 LangGraph deps",
         "  • Added strands-agents>=1.29.0",
         "  • Removed langgraph-cli (dev)",
     ], head_color=YELLOW, line_size=12)

# CREATED
card(s, 9.3, 1.2, 3.7, 4.8, "✨  CREATED",
     [
         "ee_agent_strands.py",
         "  ~520 lines",
         "",
         "  Classes:",
         "  • ManuscriptSubmission",
         "  • EditorAssignmentAgent",
         "",
         "  Methods:",
         "  • async_execute_workflow()",
         "  • _build_tools()",
         "  • _build_system_prompt()",
         "  • _build_session_manager()",
         "  • _save_to_long_term_memory()",
         "  • _call_assign_api()",
     ], head_color=GREEN, line_size=12)

# pyproject diff footer
box(s, 0.4, 6.1, 12.53, 0.95, fill_color=RGBColor(0x12, 0x1A, 0x2E),
    line_color=MID_BLUE, line_width=Pt(1))
txt(s, "pyproject.toml deps:  "
       "REMOVED → langgraph, langchain-aws, langchain-mcp-adapters, anthropic, langgraph-cli  |  "
       "KEPT → langgraph-checkpoint-postgres (for AsyncPostgresStore)  |  "
       "ADDED → strands-agents>=1.29.0",
    0.6, 6.15, 12.2, 0.5, size=11, color=LIGHT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Live Data & Learning Loop Proof
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Live Data — Learning Loop Proof", "Data from POC-memoryImplementation runs (Feb–Mar 2026)")

# Timeline boxes
def timeline_box(slide, l, t, w, h, date, event, detail, status_color):
    box(slide, l, t, w, h, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
        line_color=status_color, line_width=Pt(2))
    txt(slide, date, l+0.1, t+0.1, w-0.2, 0.3, size=12, bold=True, color=status_color)
    txt(slide, event, l+0.1, t+0.4, w-0.2, 0.35, size=13, bold=True, color=WHITE)
    txt(slide, detail, l+0.1, t+0.75, w-0.2, h-0.85,
        size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

timeline_box(s, 0.4, 1.2, 3.7, 2.8,
    "Feb 23–24, 2026", "Memory System Built",
    "• AsyncPostgresSaver (Tier 2)\n• AsyncPostgresStore (Tier 3)\n"
    "• 6 tables created in mspubs schema\n• 6 unit + 4 integration tests pass",
    LIGHT_BLUE)

timeline_box(s, 4.5, 1.2, 4.3, 2.8,
    "Feb 24, 2026", "First Run → Store Populated",
    "• v1.2.6 deployed with memory\n• First workflow run completed\n"
    "• store table: 5 rows saved\n• checkpoint table: 48 rows",
    YELLOW)

timeline_box(s, 9.2, 1.2, 3.73, 2.8,
    "Feb 26, 2026", "🎯 Learning Loop PROVEN",
    "• Pod log: 'Long-term memory search\n  returned 1 results → injecting\n"
    "  1 similar past assignments'\n• AI used its own history!",
    GREEN)

# DB snapshot table
box(s, 0.4, 4.2, 12.53, 2.85, fill_color=RGBColor(0x12, 0x1A, 0x2E),
    line_color=MID_BLUE, line_width=Pt(1))
txt(s, "Mar 2, 2026 Database Snapshot (mspubs-dev RDS)", 0.6, 4.25,
    12.0, 0.35, size=14, bold=True, color=LIGHT_BLUE)

headers = ["Table", "Rows", "Purpose", "Status"]
col_x2 = [0.6, 3.5, 5.0, 10.8]
col_w2 = [2.7, 1.3, 5.5, 2.0]

for ci, (h2, cx, cw) in enumerate(zip(headers, col_x2, col_w2)):
    txt(s, h2, cx, 4.65, cw, 0.28, size=12, bold=True, color=AWS_ORANGE)

rows = [
    ("mspubs.checkpoint",           "48",  "LangGraph session checkpoints",  "Retired (Strands uses S3)"),
    ("mspubs.checkpoint_blobs",     "96",  "LangGraph blob storage",         "Retired"),
    ("mspubs.checkpoint_writes",    "0",   "Write tracking",                 "Retired"),
    ("mspubs.store",                "5",   "Long-term assignment memory",    "✅  ACTIVE — Strands uses this"),
    ("mspubs.store_migrations",     "1",   "Schema version control",         "Active"),
]
y = 4.95
for row in rows:
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x2, col_w2)):
        cell_color = GREEN if "ACTIVE" in cell else (
            RGBColor(0x88, 0x88, 0x88) if "Retired" in cell else WHITE)
        txt(s, cell, cx, y, cw, 0.26, size=11, color=cell_color)
    y += 0.27


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Session Memory Workaround
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Session Memory Workaround (No Infra Required)",
             "Strands degrades gracefully — S3 is optional")

# Problem
box(s, 0.4, 1.2, 5.9, 2.6, fill_color=RGBColor(0x2A, 0x1A, 0x1A),
    line_color=RED, line_width=Pt(1.5))
txt(s, "❌  Blocker", 0.6, 1.25, 5.5, 0.35, size=14, bold=True, color=RED)
for y, line in zip([1.65, 1.95, 2.25, 2.55, 2.85],
                   ["IRSA role (er-bedrock-sa) needs s3:PutObject / s3:GetObject",
                    "on acs-us-east-1-gtsai-dev-s3-er-sessions",
                    "Infra team not available to add the policy",
                    "Without permission: S3SessionManager raises AccessDenied",
                    "Result: pod crashes or no session memory"]):
    txt(s, line, 0.6, y, 5.6, 0.3, size=12, color=WHITE)

# Arrow
txt(s, "→", 6.5, 2.3, 0.7, 0.5, size=28, bold=True,
    color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Solution
box(s, 7.2, 1.2, 5.73, 2.6, fill_color=RGBColor(0x12, 0x2A, 0x1A),
    line_color=GREEN, line_width=Pt(1.5))
txt(s, "✅  Built-in Graceful Degradation", 7.4, 1.25, 5.4, 0.35,
    size=14, bold=True, color=GREEN)
for y, line in zip([1.65, 1.95, 2.25, 2.55, 2.85],
                   ["S3_SESSIONS_BUCKET env var not set → session manager = None",
                    "Agent runs without session persistence",
                    "Tier 3 (Postgres) long-term memory still fully works",
                    "No code changes needed — already built this way",
                    "Just leave S3_SESSIONS_BUCKET unset in ConfigMap"]):
    txt(s, line, 7.4, y, 5.4, 0.3, size=12, color=WHITE)

# Code block
box(s, 0.4, 4.0, 12.53, 2.1, fill_color=RGBColor(0x0D, 0x11, 0x17),
    line_color=LIGHT_BLUE, line_width=Pt(1))
txt(s, "ee_agent_strands.py — _build_session_manager()", 0.6, 4.05,
    12.0, 0.3, size=12, bold=True, color=LIGHT_BLUE)

code_lines = [
    ("def _build_session_manager(self, session_id: str):", WHITE),
    ("    s3_bucket = os.getenv('S3_SESSIONS_BUCKET')   # ← not set = None", LIGHT_BLUE),
    ("    if not s3_bucket:", WHITE),
    ("        return None   # ← graceful: agent runs without session memory", GREEN),
    ("    ...               # ← if set: S3SessionManager activated", YELLOW),
]
y = 4.40
for code, color in code_lines:
    txt(s, code, 0.7, y, 12.0, 0.25, size=11, color=color)
    y += 0.26

# Recommendation box
box(s, 0.4, 6.2, 12.53, 0.85, fill_color=RGBColor(0x1A, 0x2A, 0x1A),
    line_color=GREEN, line_width=Pt(1))
txt(s,
    "Recommendation: Deploy without S3_SESSIONS_BUCKET for now. "
    "Tier 3 learning loop remains fully active. "
    "Add S3 permission later when infra team is available.",
    0.6, 6.25, 12.2, 0.7, size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Advantages Comparison Table
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "LangGraph vs Strands — Detailed Comparison")

headers = ["Dimension", "LangGraph (v1.3.0)", "Strands (v1.4.0)", "Winner"]
col_x3 = [0.3, 3.0, 7.5, 11.7]
col_w3 = [2.6, 4.3, 4.0, 1.55]
y = 1.15

# Header row
box(s, 0.3, 1.1, 12.73, 0.38, fill_color=MID_BLUE)
for h3, cx, cw in zip(headers, col_x3, col_w3):
    txt(s, h3, cx+0.05, 1.13, cw, 0.32, size=13, bold=True, color=WHITE)

comparison = [
    ("Code volume",        "3 files, ~1,094 lines",          "1 file, ~520 lines",                "Strands"),
    ("Workflow type",      "Graph: fixed node sequence",      "Model-driven: LLM decides",         "Strands"),
    ("Tier 2 storage",     "Postgres (6 RDS tables)",         "S3 (1 JSON/run)",                   "Strands"),
    ("Tier 3 storage",     "AsyncPostgresStore (RDS)",        "AsyncPostgresStore (RDS)",          "Tie"),
    ("AWS integration",    "Anthropic SDK + langchain-aws",   "Native boto3/IRSA only",            "Strands"),
    ("Dependency count",   "+5 packages (langgraph stack)",   "+1 package (strands-agents)",       "Strands"),
    ("Tool flexibility",   "Fixed tool call order",           "LLM skips unnecessary tools",       "Strands"),
    ("Learning loop",      "Proven Feb 24→26",                "Preserved (same Tier 3)",           "Tie"),
    ("Pod restart",        "Loses session state",             "S3 reloads session (when enabled)", "Strands"),
    ("Pydantic compat.",   "V1/V2 mix — warnings",            "V2 ConfigDict — clean",             "Strands"),
    ("Test coverage",      "6 unit + 4 integration",         "Inherited + new unit tests TBD",    "LangGraph"),
    ("Maturity",           "Stable, widely used",             "v1.29.0, AWS-backed, 5.3k stars",  "LangGraph"),
]

y = 1.52
for i, row in enumerate(comparison):
    row_bg = RGBColor(0x1E, 0x2D, 0x4A) if i % 2 == 0 else RGBColor(0x16, 0x22, 0x38)
    box(s, 0.3, y, 12.73, 0.34, fill_color=row_bg)
    for cell, cx, cw in zip(row, col_x3, col_w3):
        color = (GREEN if cell == "Strands" else
                 RED if cell == "LangGraph" else
                 YELLOW if cell == "Tie" else WHITE)
        bold = cell in ("Strands", "LangGraph", "Tie")
        txt(s, cell, cx+0.05, y+0.03, cw, 0.28, size=11,
            bold=bold, color=color)
    y += 0.35


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Dev Deployment Checklist
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Dev Deployment Checklist", "EKS namespace: er  |  image: ghcr.io/acspubsedsg/at-ai-editor-recommender")

items_left = [
    ("✅", "Branch pushed",               "feature/strands-memory-implementation @ 5c927636"),
    ("✅", "ee_graph_*.py deleted",        "ee_graph_anthropic, ee_graph_claude, ee_graph all removed"),
    ("✅", "ee_agent_strands.py created",  "EditorAssignmentAgent with @tool, BedrockModel, S3SessionManager"),
    ("✅", "LangGraph deps removed",       "langgraph, langchain-aws, anthropic removed from pyproject.toml"),
    ("✅", "strands-agents added",         "strands-agents>=1.29.0 in pyproject.toml"),
    ("✅", "Pydantic V2 fixed",            "ConfigDict(from_attributes=True) — zero warnings"),
    ("✅", "Import check passes",          "warnings.filterwarnings('error') → ALL IMPORTS CLEAN"),
    ("✅", "S3 bucket created",            "acs-us-east-1-gtsai-dev-s3-er-sessions (us-east-1)"),
    ("✅", "K8s YAML updated",             "S3_SESSIONS_BUCKET set in deployment-irsa.yaml"),
]

items_right = [
    ("⬜", "Build new Docker image",       "Bump version v1.3.0 → v1.4.0 in Dockerfile/pyproject.toml"),
    ("⬜", "Push image to GHCR",           "ghcr.io/acspubsedsg/at-ai-editor-recommender:v1.4.0"),
    ("⬜", "Update image tag in YAML",     "deployment-irsa.yaml: image: ...v1.4.0"),
    ("⬜", "kubectl apply",                "kubectl apply -f k8s/dev/at-ai-editor-recommender/"),
    ("⬜", "Verify pod starts",            "kubectl get pods -n er — check Running status"),
    ("⬜", "Check logs",                   "kubectl logs -n er <pod> — look for 'Strands workflow starting'"),
    ("⬜", "Run test workflow",            "POST /execute_workflow with a test manuscript"),
    ("⬜", "Add IRSA S3 policy (infra)",   "s3:PutObject/GetObject on acs-us-east-1-gtsai-dev-s3-er-sessions"),
    ("⬜", "Verify S3 session saved",      "Check S3 bucket for session JSON after first run"),
]

y = 1.25
for icon, title, detail in items_left:
    color = GREEN if icon == "✅" else YELLOW
    txt(s, f"{icon} {title}", 0.4, y, 4.5, 0.28, size=12, bold=True, color=color)
    txt(s, detail, 4.9, y, 3.5, 0.28, size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.33

y = 1.25
for icon, title, detail in items_right:
    color = YELLOW
    txt(s, f"{icon} {title}", 6.9, y, 3.5, 0.28, size=12, bold=True, color=color)
    txt(s, detail, 10.4, y, 2.7, 0.28, size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.33


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_NAVY)
slide_header(s, "Next Steps")

steps = [
    ("🚀", "Immediate",
     "Deploy to dev EKS (v1.4.0)",
     "Build image, push to GHCR, kubectl apply — no IRSA S3 needed yet"),
    ("🧪", "Immediate",
     "Run test workflow & verify Tier 3 memory",
     "POST /execute_workflow → check logs for 'Strands workflow starting' + Postgres store write"),
    ("🔐", "When infra available",
     "Add IRSA S3 policy for session memory",
     "s3:PutObject + s3:GetObject on acs-us-east-1-gtsai-dev-s3-er-sessions"),
    ("📊", "After 5+ runs",
     "Prove Strands learning loop (same as LangGraph proof)",
     "Check pod logs for: 'Long-term memory search returned N results → injecting N similar past assignments'"),
    ("🔄", "Stage / Prod",
     "Promote branch → PR → merge → deploy to stage",
     "Update stage K8s ConfigMap with POSTGRES_URI, S3_SESSIONS_BUCKET, MODEL_ID"),
    ("🧹", "Cleanup",
     "Remove retired Postgres checkpoint tables",
     "DROP checkpoint, checkpoint_blobs, checkpoint_writes tables from mspubs schema"),
    ("📝", "Documentation",
     "Update README with Strands architecture diagram",
     "Add IRSA policy snippet + S3 bucket requirements for new environments"),
]

y = 1.2
for icon, timing, title, detail in steps:
    timing_color = (GREEN if timing == "Immediate" else
                    YELLOW if "infra" in timing else
                    LIGHT_BLUE)
    box(s, 0.4, y, 1.2, 0.55, fill_color=RGBColor(0x1E, 0x2D, 0x4A),
        line_color=timing_color, line_width=Pt(1))
    txt(s, icon, 0.4, y, 1.2, 0.28, size=18, align=PP_ALIGN.CENTER, color=WHITE)
    txt(s, timing, 0.4, y+0.28, 1.2, 0.25, size=9, bold=True,
        color=timing_color, align=PP_ALIGN.CENTER)
    txt(s, title, 1.75, y+0.02, 4.5, 0.28, size=13, bold=True, color=WHITE)
    txt(s, detail, 1.75, y+0.30, 10.8, 0.28, size=11,
        color=RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.70


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = r"c:\Users\szs17\OneDrive - ACS\Documents\GitHub\at-ai-editor-recommender\docs\Strands_Migration_Editor_Recommender.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
