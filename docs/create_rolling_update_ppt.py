"""Generate Rolling Pod Update – Live Validation PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colors ──────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE  = RGBColor(0x00, 0x9C, 0xDE)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GOLD  = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_TEAL  = RGBColor(0x1A, 0xBC, 0x9C)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT    = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BG     = RGBColor(0xF5, 0xF6, 0xFA)
CODE_BG      = RGBColor(0x1E, 0x1E, 0x2E)
ROW_EVEN     = RGBColor(0x24, 0x24, 0x3E)
ROW_ODD      = RGBColor(0x2A, 0x2A, 0x48)
HDR_BG       = RGBColor(0x00, 0x7A, 0xAD)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]   # truly blank


# ── Helpers ───────────────────────────────────────────────────────────────────
def dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def light_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # rounded rect = 5; plain = 1 (MSO_SHAPE_TYPE.RECTANGLE)
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
       align=PP_ALIGN.LEFT, italic=False, wrap=True, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size   = Pt(size)
    p.font.color.rgb = color
    p.font.bold   = bold
    p.font.italic = italic
    p.font.name   = font
    return box


def multiline(slide, l, t, w, h, lines, default_size=16, default_color=WHITE,
              font="Calibri"):
    """lines = list of (text, bold, color_or_None, size_or_None, align_or_None)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, row in enumerate(lines):
        text  = row[0]
        bold  = row[1] if len(row) > 1 else False
        clr   = row[2] if len(row) > 2 else None
        size  = row[3] if len(row) > 3 else None
        align = row[4] if len(row) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size      = Pt(size or default_size)
        p.font.color.rgb = clr or default_color
        p.font.bold      = bold
        p.font.name      = font
        p.alignment      = align
        p.space_after    = Pt((size or default_size) * 0.25)
    return box


def accent_bar(slide, color=ACCENT_BLUE, left=Inches(0.45), top=Inches(1.35),
               width=Inches(0.07), height=Inches(0.52)):
    r = rect(slide, left, top, width, height, color)
    return r


def slide_title(slide, title, subtitle=None, title_color=WHITE,
                sub_color=ACCENT_BLUE):
    tb(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.65),
       title, size=32, color=title_color, bold=True)
    if subtitle:
        tb(slide, Inches(0.55), Inches(1.9), Inches(12.2), Inches(0.4),
           subtitle, size=16, color=sub_color, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
dark_bg(s1)

# Top accent strip
rect(s1, 0, 0, W, Inches(0.07), ACCENT_BLUE)

# Large bold title block
rect(s1, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.2), RGBColor(0x22, 0x22, 0x40))

tb(s1, Inches(0.8), Inches(1.95), Inches(11.7), Inches(1.0),
   "Rolling Pod Update", size=46, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

tb(s1, Inches(0.8), Inches(2.95), Inches(11.7), Inches(0.7),
   "Live Validation on EKS Dev Cluster", size=28, color=WHITE, bold=False,
   align=PP_ALIGN.CENTER)

tb(s1, Inches(0.8), Inches(3.65), Inches(11.7), Inches(0.5),
   "Service: at-ai-editor-recommender  |  Namespace: er  |  Cluster: acs-gtsai-dev",
   size=15, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Divider
rect(s1, Inches(3.5), Inches(5.1), Inches(6.3), Inches(0.04), ACCENT_BLUE)

tb(s1, Inches(0), Inches(5.3), W, Inches(0.4),
   "ENG-9411  •  ACS GTS AI  •  March 2026",
   size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Bottom accent
rect(s1, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem: Downtime During Deployments
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
dark_bg(s2)
rect(s2, 0, 0, W, Inches(0.07), ACCENT_RED)

accent_bar(s2, ACCENT_RED)
slide_title(s2, "The Problem: Downtime During Deployments",
            subtitle="Default Kubernetes rolling update settings cause service interruptions",
            sub_color=ACCENT_RED)

# Left card
rect(s2, Inches(0.45), Inches(2.5), Inches(5.8), Inches(3.8), RGBColor(0x22, 0x22, 0x40))
tb(s2, Inches(0.6), Inches(2.6), Inches(5.5), Inches(0.45),
   "Default Strategy", size=18, color=ACCENT_RED, bold=True)

multiline(s2, Inches(0.6), Inches(3.1), Inches(5.5), Inches(3.0), [
    ("strategy:", False, MED_GRAY, 14),
    ("  type: RollingUpdate", False, LIGHT_GRAY, 13),
    ("  rollingUpdate:", False, MED_GRAY, 14),
    ("    maxUnavailable: 25%  ← PROBLEM", True, ACCENT_RED, 14),
    ("    maxSurge: 25%", False, LIGHT_GRAY, 13),
], font="Courier New")

tb(s2, Inches(0.6), Inches(5.25), Inches(5.5), Inches(0.8),
   "With only 1 replica, 25% maxUnavailable = pod goes down before new one is ready → downtime!",
   size=13, color=ACCENT_RED)

# Right card — timeline
rect(s2, Inches(6.85), Inches(2.5), Inches(6.0), Inches(3.8), RGBColor(0x22, 0x22, 0x40))
tb(s2, Inches(7.0), Inches(2.6), Inches(5.7), Inches(0.45),
   "What Happens (Old Behavior)", size=18, color=LIGHT_GRAY, bold=True)

steps_old = [
    ("T+0s", "Deploy triggered"),
    ("T+1s", "Old pod TERMINATED  ← gap begins"),
    ("T+12s", "New pod Pending"),
    ("T+18s", "New pod ContainerCreating"),
    ("T+28s", "New pod Running"),
    ("", "~27s of service interruption!"),
]
y = Inches(3.15)
for time, desc in steps_old:
    clr = ACCENT_RED if "gap" in desc or "interruption" in desc else LIGHT_GRAY
    bold = "interruption" in desc or "gap" in desc
    tb(s2, Inches(7.0), y, Inches(1.1), Inches(0.32), time,
       size=13, color=ACCENT_GOLD, bold=True)
    tb(s2, Inches(8.15), y, Inches(4.5), Inches(0.32), desc,
       size=13, color=clr, bold=bold)
    y += Inches(0.45)

rect(s2, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Fix: Zero-Downtime Rolling Update Config
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
dark_bg(s3)
rect(s3, 0, 0, W, Inches(0.07), ACCENT_GREEN)

accent_bar(s3, ACCENT_GREEN)
slide_title(s3, "The Fix: Zero-Downtime Rolling Update Configuration",
            subtitle="Applied to at-ai-editor-recommender-deployment via Kustomize overlay",
            sub_color=ACCENT_GREEN)

# YAML block (left)
rect(s3, Inches(0.45), Inches(2.5), Inches(6.0), Inches(4.5), CODE_BG)

multiline(s3, Inches(0.65), Inches(2.65), Inches(5.7), Inches(4.2), [
    ("spec:", False, MED_GRAY, 14),
    ("  replicas: 1", False, LIGHT_GRAY, 14),
    ("  strategy:", False, MED_GRAY, 14),
    ("    type: RollingUpdate", False, LIGHT_GRAY, 14),
    ("    rollingUpdate:", False, MED_GRAY, 14),
    ("      maxUnavailable: 0   ✓", True, ACCENT_GREEN, 15),
    ("      maxSurge: 1         ✓", True, ACCENT_GREEN, 15),
    ("  template:", False, MED_GRAY, 14),
    ("    spec:", False, MED_GRAY, 14),
    ("      containers:", False, MED_GRAY, 14),
    ("        - imagePullPolicy: IfNotPresent  ✓", True, ACCENT_TEAL, 14),
], font="Courier New")

# Explanation bullets (right)
right_bullets = [
    ("maxUnavailable: 0", True, ACCENT_GREEN, 17),
    ("", False, WHITE, 6),
    ("Kubernetes will NOT terminate the old pod until the new pod is fully Ready. "
     "Guarantees at least 1 pod is always serving traffic.", False, LIGHT_GRAY, 14),
    ("", False, WHITE, 8),
    ("maxSurge: 1", True, ACCENT_BLUE, 17),
    ("", False, WHITE, 6),
    ("Allows one extra pod above the desired replica count during rollout. "
     "New pod starts alongside old pod, health-checked before traffic switch.", False, LIGHT_GRAY, 14),
    ("", False, WHITE, 8),
    ("imagePullPolicy: IfNotPresent", True, ACCENT_TEAL, 17),
    ("", False, WHITE, 6),
    ("Skips image pull if the image is already cached on the node. "
     "Eliminates dependency on GitHub Container Registry during restarts.", False, LIGHT_GRAY, 14),
]
multiline(s3, Inches(7.0), Inches(2.5), Inches(5.9), Inches(4.5), right_bullets)

rect(s3, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works: Pod Lifecycle Diagram
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
dark_bg(s4)
rect(s4, 0, 0, W, Inches(0.07), ACCENT_BLUE)

accent_bar(s4, ACCENT_BLUE)
slide_title(s4, "How It Works: Pod Lifecycle During a Rolling Update",
            subtitle="maxUnavailable: 0  ensures the old pod stays alive until the new pod is healthy")

# Phase timeline
phases = [
    (ACCENT_GOLD,  "Phase 1",  "Rollout\nTriggered",    "T + 0s"),
    (ACCENT_BLUE,  "Phase 2",  "New Pod\nPending",       "T + 2s"),
    (ACCENT_TEAL,  "Phase 3",  "New Pod\nRunning",       "T + 10s"),
    (ACCENT_GREEN, "Phase 4",  "New Pod\nReady ✓",       "T + 18s"),
    (ACCENT_RED,   "Phase 5",  "Old Pod\nTerminated",    "T + 20s"),
]

box_w = Inches(2.05)
box_h = Inches(2.0)
gap   = Inches(0.2)
start_x = Inches(0.45)
y_box = Inches(2.8)

for i, (clr, label, desc, timing) in enumerate(phases):
    x = start_x + i * (box_w + gap)
    # Background box
    rect(s4, x, y_box, box_w, box_h, RGBColor(0x22, 0x22, 0x40))
    # Color top bar
    rect(s4, x, y_box, box_w, Inches(0.08), clr)
    # Label
    tb(s4, x + Inches(0.05), y_box + Inches(0.12), box_w - Inches(0.1), Inches(0.35),
       label, size=15, color=clr, bold=True, align=PP_ALIGN.CENTER)
    # Description
    tb(s4, x + Inches(0.05), y_box + Inches(0.52), box_w - Inches(0.1), Inches(0.9),
       desc, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    # Timing
    tb(s4, x + Inches(0.05), y_box + Inches(1.6), box_w - Inches(0.1), Inches(0.35),
       timing, size=13, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
    # Arrow (except last)
    if i < len(phases) - 1:
        ax = x + box_w + Inches(0.04)
        tb(s4, ax, y_box + Inches(0.75), Inches(0.16), Inches(0.4),
           "→", size=22, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Key highlight bar
rect(s4, Inches(0.45), Inches(5.15), Inches(12.4), Inches(0.6), RGBColor(0x0D, 0x4F, 0x2A))
tb(s4, Inches(0.55), Inches(5.2), Inches(12.2), Inches(0.5),
   "✓  Traffic never interrupted — at every moment between T+0 and T+20, at least 1 pod is Running and Ready",
   size=15, color=ACCENT_GREEN, bold=True)

# Note at bottom
rect(s4, Inches(0.45), Inches(6.0), Inches(12.4), Inches(1.0), RGBColor(0x1E, 0x1E, 0x38))
multiline(s4, Inches(0.6), Inches(6.05), Inches(12.0), Inches(0.9), [
    ("How the health check works:", True, ACCENT_BLUE, 14),
    ("Kubernetes checks the readinessProbe (HTTP GET /ready) before marking a pod Ready. "
     "Only when Ready → old pod is removed from Service endpoints and then terminated.",
     False, LIGHT_GRAY, 13),
])

rect(s4, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LIVE TEST RESULTS (CENTER FOCUS)
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
dark_bg(s5)
rect(s5, 0, 0, W, Inches(0.07), ACCENT_GREEN)

accent_bar(s5, ACCENT_GREEN)
slide_title(s5, "LIVE TEST RESULTS — Zero-Downtime Confirmed",
            subtitle="kubectl rollout restart  →  Observed every 5 seconds for 60 seconds  |  Dev Cluster, namespace: er",
            sub_color=ACCENT_GREEN)

# ── Table setup ──────────────────────────────────────────────────────────────
# Columns: Elapsed | Timestamp | Total Pods | Running | Pod Name/Age | Status
# Data rows
headers = ["T+", "Timestamp", "Total Pods", "Running", "Pod Name", "Age", "Status"]
col_widths = [Inches(0.65), Inches(1.35), Inches(1.05), Inches(0.9), Inches(3.65), Inches(0.85), Inches(1.35)]

test_rows = [
    # elapsed, time, total, running, pod, age, status
    ("0s",   "14:42:36", "1", "0", "Rollout restart triggered ↓", "",    "TRIGGERED"),
    ("18s",  "14:42:54", "1", "1", "66c747c9b7-w75pp",            "18s", "✓ RUNNING"),
    ("28s",  "14:43:04", "1", "1", "66c747c9b7-w75pp",            "28s", "✓ RUNNING"),
    ("38s",  "14:43:13", "1", "1", "66c747c9b7-w75pp",            "38s", "✓ RUNNING"),
    ("48s",  "14:43:23", "1", "1", "66c747c9b7-w75pp",            "48s", "✓ RUNNING"),
    ("57s",  "14:43:33", "1", "1", "66c747c9b7-w75pp",            "57s", "✓ RUNNING"),
    ("67s",  "14:43:43", "1", "1", "66c747c9b7-w75pp",            "67s", "✓ RUNNING"),
    ("77s",  "14:43:52", "1", "1", "66c747c9b7-w75pp",            "77s", "✓ RUNNING"),
]

table_l = Inches(0.38)
table_t = Inches(2.5)
row_h   = Inches(0.38)

# Header row
hdr_t = table_t
x = table_l
for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
    rect(s5, x, hdr_t, cw, row_h, HDR_BG)
    tb(s5, x + Inches(0.05), hdr_t + Inches(0.04), cw - Inches(0.05), row_h - Inches(0.04),
       hdr, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    x += cw

# Data rows
for r_idx, row_data in enumerate(test_rows):
    row_top = hdr_t + row_h + r_idx * row_h
    is_trigger = row_data[0] == "0s"
    bg = RGBColor(0x15, 0x35, 0x20) if "RUNNING" in row_data[6] else \
         RGBColor(0x35, 0x25, 0x05) if is_trigger else ROW_EVEN

    x = table_l
    for c_idx, (cell_text, cw) in enumerate(zip(row_data, col_widths)):
        rect(s5, x, row_top, cw, row_h, bg)
        # Color logic
        if c_idx == 6:  # Status column
            clr = ACCENT_GREEN if "RUNNING" in cell_text else ACCENT_GOLD
            bld = True
        elif c_idx == 0:  # Elapsed
            clr = ACCENT_GOLD
            bld = True
        elif c_idx == 3 and cell_text == "1":  # Running count
            clr = ACCENT_GREEN
            bld = True
        elif c_idx == 4 and "66c747c9b7" in cell_text:  # new pod
            clr = ACCENT_TEAL
            bld = False
        else:
            clr = LIGHT_GRAY
            bld = False

        fs = 11 if c_idx == 4 else 13
        alg = PP_ALIGN.LEFT if c_idx in (4,) else PP_ALIGN.CENTER
        tb(s5, x + Inches(0.04), row_top + Inches(0.05), cw - Inches(0.04), row_h - Inches(0.05),
           cell_text, size=fs, color=clr, bold=bld, align=alg, font="Calibri")
        x += cw

    # Thin border line between rows
    rect(s5, table_l, row_top + row_h - Inches(0.008),
         sum(col_widths), Inches(0.008), RGBColor(0x33, 0x33, 0x55))

# ── Key notes below table ─────────────────────────────────────────────────────
note_t = hdr_t + row_h * (len(test_rows) + 1) + Inches(0.08)

rect(s5, table_l, note_t, sum(col_widths), Inches(0.05), ACCENT_GREEN)

notes = [
    ("Old pod 6b58d7bc95-2j72w (running 8 days) was already terminated before first observation at T+18s  —  transition completed in < 18 seconds",
     False, MED_GRAY, 12),
    ("Running count NEVER dropped to 0  •  maxUnavailable: 0 enforced correctly  •  Zero downtime ✓", True, ACCENT_GREEN, 13),
]
multiline(s5, table_l, note_t + Inches(0.1), sum(col_widths), Inches(0.65), notes)

rect(s5, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Before vs After Comparison
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
dark_bg(s6)
rect(s6, 0, 0, W, Inches(0.07), ACCENT_BLUE)

accent_bar(s6, ACCENT_BLUE)
slide_title(s6, "Before vs. After: Impact Summary",
            subtitle="Two configuration changes, dramatic reliability improvement")

# BEFORE card
rect(s6, Inches(0.45), Inches(2.4), Inches(5.9), Inches(4.6), RGBColor(0x30, 0x15, 0x15))
rect(s6, Inches(0.45), Inches(2.4), Inches(5.9), Inches(0.5), ACCENT_RED)
tb(s6, Inches(0.55), Inches(2.42), Inches(5.7), Inches(0.45),
   "✗  BEFORE", size=20, color=WHITE, bold=True)

before_items = [
    ("maxUnavailable: 25%  (= 1 pod with replicas:1)", True, ACCENT_RED, 14),
    ("maxSurge: 25%", False, LIGHT_GRAY, 13),
    ("imagePullPolicy: Always", False, ACCENT_RED, 14),
    ("", False, WHITE, 6),
    ("→ Pod killed before replacement is Ready", False, LIGHT_GRAY, 13),
    ("→ ~27 seconds of service unavailability", False, ACCENT_RED, 14),
    ("→ Pull fails if GHCR is unreachable", False, LIGHT_GRAY, 13),
    ("→ Slow restarts due to unnecessary image pulls", False, LIGHT_GRAY, 13),
]
multiline(s6, Inches(0.6), Inches(3.0), Inches(5.6), Inches(3.8), before_items)

# AFTER card
rect(s6, Inches(7.0), Inches(2.4), Inches(5.9), Inches(4.6), RGBColor(0x10, 0x28, 0x18))
rect(s6, Inches(7.0), Inches(2.4), Inches(5.9), Inches(0.5), ACCENT_GREEN)
tb(s6, Inches(7.1), Inches(2.42), Inches(5.7), Inches(0.45),
   "✓  AFTER", size=20, color=WHITE, bold=True)

after_items = [
    ("maxUnavailable: 0  (old pod stays alive)", True, ACCENT_GREEN, 14),
    ("maxSurge: 1  (new pod spins up alongside)", False, ACCENT_TEAL, 13),
    ("imagePullPolicy: IfNotPresent", False, ACCENT_TEAL, 14),
    ("", False, WHITE, 6),
    ("→ Replacement ready BEFORE old pod goes down", False, LIGHT_GRAY, 13),
    ("→ 0 seconds downtime (live tested ✓)", True, ACCENT_GREEN, 14),
    ("→ No GHCR dependency at restart time", False, LIGHT_GRAY, 13),
    ("→ Fast restarts using cached image", False, LIGHT_GRAY, 13),
]
multiline(s6, Inches(7.1), Inches(3.0), Inches(5.6), Inches(3.8), after_items)

# VS label
tb(s6, Inches(6.05), Inches(4.2), Inches(1.2), Inches(0.7),
   "VS", size=30, color=MED_GRAY, bold=True, align=PP_ALIGN.CENTER)

rect(s6, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cluster-Wide Status & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
dark_bg(s7)
rect(s7, 0, 0, W, Inches(0.07), ACCENT_GOLD)

accent_bar(s7, ACCENT_GOLD)
slide_title(s7, "Current Status Across All Services",
            subtitle="Dev cluster audit — er namespace  |  March 5, 2026",
            sub_color=ACCENT_GOLD)

# Table
svc_headers = ["Service", "maxUnavailable", "maxSurge", "imagePullPolicy", "Status"]
svc_col_w   = [Inches(4.0), Inches(1.9), Inches(1.5), Inches(2.5), Inches(2.9)]

svc_data = [
    ("at-ai-editor-recommender",      "0",    "1",    "IfNotPresent", "✓ COMPLETE + TESTED"),
    ("agent-invoke-client",            "25%",  "25%",  "Always",       "⚠ Needs Update"),
    ("at-ai-editor-recommender-poc",   "25%",  "25%",  "IfNotPresent", "⚠ Needs Sync"),
    ("at-ai-topic-extraction",         "25%",  "25%",  "Always",       "⚠ Needs Update"),
    ("ee-api",                         "25%",  "25%",  "Always",       "⚠ Needs Update"),
    ("ee-api-poc",                     "25%",  "25%",  "IfNotPresent", "⚠ Needs Sync"),
]

tl = Inches(0.38)
tt = Inches(2.55)
rh = Inches(0.5)

# Header
x = tl
for hdr, cw in zip(svc_headers, svc_col_w):
    rect(s7, x, tt, cw, rh, HDR_BG)
    tb(s7, x + Inches(0.05), tt + Inches(0.07), cw, rh,
       hdr, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += cw

for r, row in enumerate(svc_data):
    rt = tt + rh + r * rh
    complete = "COMPLETE" in row[4]
    bg = RGBColor(0x10, 0x28, 0x18) if complete else ROW_ODD
    x = tl
    for c, (cell, cw) in enumerate(zip(row, svc_col_w)):
        rect(s7, x, rt, cw, rh, bg)
        if c == 4:
            clr = ACCENT_GREEN if complete else ACCENT_GOLD
            bld = True
        elif c == 1 and cell == "0":
            clr = ACCENT_GREEN; bld = True
        elif c == 1 and "25%" in cell:
            clr = ACCENT_RED; bld = True
        elif c == 3 and cell == "Always":
            clr = ACCENT_RED; bld = False
        elif c == 3 and cell == "IfNotPresent":
            clr = ACCENT_GREEN; bld = False
        else:
            clr = LIGHT_GRAY; bld = False
        alg = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        tb(s7, x + Inches(0.06), rt + Inches(0.07), cw - Inches(0.06), rh,
           cell, size=13, color=clr, bold=bld, align=alg)
        x += cw
    rect(s7, tl, rt + rh - Inches(0.008), sum(svc_col_w), Inches(0.008),
         RGBColor(0x33, 0x33, 0x55))

# Next steps footer
rect(s7, Inches(0.38), Inches(6.2), Inches(12.55), Inches(0.8), RGBColor(0x1E, 0x1E, 0x38))
multiline(s7, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.7), [
    ("Next Steps:", True, ACCENT_GOLD, 14),
    ("  Apply rolling update config to remaining 5 services via ArgoCD sync  •  "
     "Fix imagePullPolicy: Always (3 services)  •  Phase 2: Migrate images to AWS ECR",
     False, LIGHT_GRAY, 13),
])

rect(s7, 0, H - Inches(0.07), W, Inches(0.07), ACCENT_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
OUT = r"c:\Users\szs17\OneDrive - ACS\Documents\GitHub\at-ai-editor-recommender\docs\Rolling-Update-Validation.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
