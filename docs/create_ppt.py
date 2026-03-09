"""Generate the Memory Implementation POC PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)      # Dark navy
ACCENT_BLUE = RGBColor(0x00, 0x9C, 0xDE)   # Bright blue
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)  # Green for checkmarks
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)    # Red for X marks
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
CODE_BG = RGBColor(0x2D, 0x2D, 0x3F)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_dark_bg(slide):
    """Fill slide with dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_light_bg(slide):
    """Fill slide with light background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=WHITE, font_name="Calibri", line_spacing=1.2):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr or color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.3)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Editor Recommender", font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Session & Long-term Memory Implementation", font_size=28, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "POC — Proof of Concept", font_size=22, color=ORANGE, bold=True)

lines = [
    ("LangGraph Checkpointer  •  LangGraph Store  •  pgvector  •  PostgreSQL", False, LIGHT_GRAY),
    ("", False, None),
    ("February 23, 2026  •  Validated on MSPUBS Dev RDS", False, MED_GRAY),
]
add_multiline_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(1.5), lines, font_size=18)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2: Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The Problem — Before Memory", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "The workflow only had in-context state — data lived in RAM during a single run",
             font_size=16, color=LIGHT_GRAY)

# Three problem boxes
problems = [
    ("❌  No Crash Recovery", "If pod restarts mid-workflow,\nall progress is lost.\nMust re-run from scratch."),
    ("❌  No Audit Trail", "No record of what happened\nat each decision step.\nCannot debug or comply."),
    ("❌  No Learning", "System cannot recall past\neditor assignments.\nEvery decision is from zero."),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.2)
    w = Inches(3.6)
    h = Inches(3.5)
    add_rounded_rect(slide, x, y, w, h, RGBColor(0x2A, 0x2A, 0x44))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.6),
                 title, font_size=20, color=ACCENT_RED, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), w - Inches(0.6), Inches(2.0),
                 desc, font_size=16, color=LIGHT_GRAY)

# Workflow bar at bottom
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
             "Current workflow:  check_resubmission  →  fetch_data  →  generate_recommendation  →  verify  →  execute_assignment",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3: Solution — Three Memory Tiers
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The Solution — Three Memory Tiers", font_size=32, color=WHITE, bold=True)

tiers = [
    ("TIER 1", "In-Context Memory", "LangGraph State TypedDict",
     "• Passes data between nodes\n• Lives in process RAM\n• Already existed",
     MED_GRAY, "EXISTED"),
    ("TIER 2", "Session Memory", "AsyncPostgresSaver → Postgres",
     "• Auto-checkpoint after every node\n• Crash recovery — resume from last step\n• Full audit trail of all state changes\n• thread_id per manuscript",
     ACCENT_BLUE, "NEW ✅"),
    ("TIER 3", "Long-term Memory", "AsyncPostgresStore → pgvector",
     "• Stores completed assignments permanently\n• Semantic search over past decisions\n• Injects past assignments into LLM prompt\n• Agent learns from history\n• Namespaced by journal",
     ACCENT_GREEN, "NEW ✅"),
]

for i, (tier, name, backend, bullets, color, badge) in enumerate(tiers):
    x = Inches(0.8 + i * 4.0)
    y = Inches(1.5)
    w = Inches(3.6)
    h = Inches(5.0)

    add_rounded_rect(slide, x, y, w, h, RGBColor(0x2A, 0x2A, 0x44))

    # Tier badge
    badge_color = ACCENT_GREEN if "NEW" in badge else MED_GRAY
    add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(1.6), Inches(0.4),
                     badge_color, badge, font_size=12, font_color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), w - Inches(0.6), Inches(0.4),
                 tier, font_size=14, color=color, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), w - Inches(0.6), Inches(0.5),
                 name, font_size=20, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.8), w - Inches(0.6), Inches(0.4),
                 backend, font_size=12, color=ORANGE)

    add_text_box(slide, x + Inches(0.3), y + Inches(2.4), w - Inches(0.6), Inches(2.2),
                 bullets, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4: Architecture Diagram
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Architecture — How Memory Integrates", font_size=32, color=WHITE, bold=True)

# FastAPI box
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.2),
                 RGBColor(0x2A, 0x2A, 0x44),
                 "FastAPI  (app.py)  —  POST /execute_workflow", font_size=16, font_color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "lifespan:  if POSTGRES_URI  →  create_checkpointer() + create_store()   else  →  run without memory (graceful)",
             font_size=12, color=LIGHT_GRAY)

# Graph nodes
nodes = [
    ("check_\nresubmission", Inches(0.5)),
    ("fetch_\nmanuscript", Inches(2.8)),
    ("generate_\nrecommendation", Inches(5.1)),
    ("verify_\nrecommendation", Inches(7.4)),
    ("execute_\nassignment", Inches(9.7)),
]

for name, x in nodes:
    add_rounded_rect(slide, x, Inches(3.0), Inches(2.0), Inches(1.0),
                     RGBColor(0x34, 0x49, 0x5E), name, font_size=13, font_color=WHITE)

# Arrows between nodes
for i in range(len(nodes) - 1):
    x1 = nodes[i][1] + Inches(2.0)
    x2 = nodes[i + 1][1]
    mid = (x1 + x2) // 2
    add_arrow(slide, x1, Inches(3.3), x2 - x1, Inches(0.3))

# Checkpoint annotation
add_text_box(slide, Inches(0.5), Inches(4.2), Inches(11.5), Inches(0.5),
             "↓ CHECKPOINT saved to Postgres after EACH node  (Tier 2: Session Memory)",
             font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Two memory boxes at bottom
# Tier 2
add_rounded_rect(slide, Inches(0.5), Inches(5.0), Inches(5.8), Inches(2.0),
                 RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(5.2), Inches(0.4),
             "TIER 2:  Session Memory  (AsyncPostgresSaver)", font_size=16, color=ACCENT_BLUE, bold=True)
lines = [
    ("• Auto-checkpoint after every node execution", False, LIGHT_GRAY),
    ("• thread_id = \"{journal_id}-{manuscript_number}\"", False, LIGHT_GRAY),
    ("• Resume from last checkpoint on pod restart", False, LIGHT_GRAY),
    ("• aget_state_history() for full audit trail", False, LIGHT_GRAY),
]
add_multiline_box(slide, Inches(0.8), Inches(5.55), Inches(5.2), Inches(1.4), lines, font_size=13)

# Tier 3
add_rounded_rect(slide, Inches(6.8), Inches(5.0), Inches(6.0), Inches(2.0),
                 RGBColor(0x1A, 0x4A, 0x2C))
add_text_box(slide, Inches(7.1), Inches(5.1), Inches(5.5), Inches(0.4),
             "TIER 3:  Long-term Memory  (AsyncPostgresStore + pgvector)", font_size=16, color=ACCENT_GREEN, bold=True)
lines = [
    ("• Saves completed assignments after workflow ends", False, LIGHT_GRAY),
    ("• Namespace: (\"assignments\", journal_id)", False, LIGHT_GRAY),
    ("• pgvector embeddings on reasoning + topics fields", False, LIGHT_GRAY),
    ("• Semantic search for similar past decisions", False, LIGHT_GRAY),
]
add_multiline_box(slide, Inches(7.1), Inches(5.55), Inches(5.5), Inches(1.4), lines, font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5: Session Memory — How Checkpointing Works
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Tier 2: Session Memory — How Checkpointing Works", font_size=32, color=ACCENT_BLUE, bold=True)

# Normal flow (left)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Normal Execution — Every Step Saved", font_size=20, color=WHITE, bold=True)

steps_normal = [
    ("① fetch_manuscript_data runs", "→ checkpoint_1 saved", ACCENT_GREEN),
    ("② generate_recommendation runs", "→ checkpoint_2 saved", ACCENT_GREEN),
    ("③ verify_recommendation runs", "→ checkpoint_3 saved", ACCENT_GREEN),
    ("④ execute_assignment runs", "→ checkpoint_4 saved", ACCENT_GREEN),
    ("⑤ DONE — full audit trail", "= 4 checkpoints", WHITE),
]

for i, (step, note, clr) in enumerate(steps_normal):
    y = Inches(1.9 + i * 0.65)
    add_text_box(slide, Inches(1.0), y, Inches(3.5), Inches(0.5),
                 step, font_size=15, color=LIGHT_GRAY)
    add_text_box(slide, Inches(4.2), y, Inches(2.5), Inches(0.5),
                 note, font_size=15, color=clr, bold=True)

# Crash recovery (right)
add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
             "Crash Recovery — Resume from Checkpoint", font_size=20, color=WHITE, bold=True)

steps_crash = [
    ("① fetch_manuscript_data runs", "→ checkpoint_1 ✅", ACCENT_GREEN),
    ("② generate_recommendation runs", "→ checkpoint_2 ✅", ACCENT_GREEN),
    ("③ Pod crashes! 💥", "💥 CRASH", ACCENT_RED),
    ("④ Pod restarts", "→ loads checkpoint_2", ORANGE),
    ("⑤ Resumes from verify_recommendation", "→ checkpoint_3 ✅", ACCENT_GREEN),
    ("⑥ execute_assignment runs", "→ checkpoint_4 ✅", ACCENT_GREEN),
    ("⑦ DONE — no work lost!", "= zero data loss", WHITE),
]

for i, (step, note, clr) in enumerate(steps_crash):
    y = Inches(1.9 + i * 0.58)
    add_text_box(slide, Inches(7.2), y, Inches(3.5), Inches(0.5),
                 step, font_size=14, color=LIGHT_GRAY)
    add_text_box(slide, Inches(10.5), y, Inches(2.5), Inches(0.5),
                 note, font_size=14, color=clr, bold=True)

# Code snippet at bottom
add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), CODE_BG)
code = ('# ee_graph_anthropic.py — async_execute_workflow()\n'
        'thread_id = f"{manuscript_submission.journal_id}-{manuscript_submission.manuscript_number}"\n'
        'config = {"configurable": {"thread_id": thread_id}}\n'
        'async for step in self._astream(manuscript_submission, config=config):  # auto-checkpoints!')
add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(1.1),
             code, font_size=13, color=ACCENT_GREEN, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6: Long-term Memory — How Assignment Storage Works
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Tier 3: Long-term Memory — Assignment Storage", font_size=32, color=ACCENT_GREEN, bold=True)

# Save flow (left)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Saving — After Workflow Completes", font_size=20, color=WHITE, bold=True)

save_items = [
    "Workflow finishes → editor assigned",
    "save_assignment_to_memory() called",
    "Extracts: editor_id, reasoning, topics",
    "Stores under namespace (\"assignments\", journal_id)",
    "Key = manuscript_number",
    "pgvector indexes reasoning + topics",
]

for i, item in enumerate(save_items):
    y = Inches(1.9 + i * 0.55)
    add_text_box(slide, Inches(1.0), y, Inches(5.5), Inches(0.5),
                 f"{'①②③④⑤⑥'[i]}  {item}", font_size=15, color=LIGHT_GRAY)

# Search flow (right)
add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
             "Searching — Future Workflows Query History", font_size=20, color=WHITE, bold=True)

search_items = [
    "New manuscript arrives",
    "search_similar_assignments(query, journal_id)",
    "pgvector finds semantically similar past decisions",
    "LLM sees: 'For similar papers, editor X worked well'",
    "Better, data-driven recommendations!",
]

for i, item in enumerate(search_items):
    y = Inches(1.9 + i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(5.5), Inches(0.5),
                 f"{'①②③④⑤'[i]}  {item}", font_size=15, color=LIGHT_GRAY)

# Data model box at bottom
add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.9), CODE_BG)
code = ('# memory.py — What gets stored per assignment\n'
        'memory_record = {\n'
        '    "editor_id": "130958",            "editor_person_id": "person-001",\n'
        '    "reasoning": "Expert in organic chemistry...",\n'
        '    "topics": "catalysis, organic synthesis",\n'
        '    "journal_id": "jacs",             "manuscript_number": "JACS-2026-00001",\n'
        '    "runner_up": "person-002",        "timestamp": "2026-02-21T10:30:00Z",\n'
        '}\n'
        'await store.aput(namespace=("assignments", "jacs"), key="JACS-2026-00001", value=memory_record)')
add_text_box(slide, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.7),
             code, font_size=13, color=ACCENT_GREEN, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6b: Memory READ — Injecting Past Assignments into LLM Prompt
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Tier 3: Memory READ — Past Assignments in LLM Prompt", font_size=32, color=ACCENT_GREEN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Before every LLM call, the system searches memory for similar past assignments and injects them into the prompt",
             font_size=16, color=LIGHT_GRAY)

# Flow steps (left)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
             "How the READ Path Works", font_size=20, color=WHITE, bold=True)

read_steps = [
    "New manuscript arrives for recommendation",
    "System extracts manuscript text (first 1000 chars)",
    "search_similar_assignments(query, journal_id, limit=5)",
    "Store returns semantically similar past assignments",
    "format_past_assignments_for_prompt(results)",
    "Formatted text injected into {past_assignments} placeholder",
    "LLM sees past decisions + makes better recommendations",
]

for i, item in enumerate(read_steps):
    y = Inches(2.4 + i * 0.5)
    add_text_box(slide, Inches(1.0), y, Inches(5.5), Inches(0.5),
                 f"{'①②③④⑤⑥⑦'[i]}  {item}", font_size=14, color=LIGHT_GRAY)

# Example prompt box (right)
add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.5),
             "What the LLM Sees in the Prompt", font_size=20, color=WHITE, bold=True)

add_rounded_rect(slide, Inches(6.8), Inches(2.4), Inches(6.0), Inches(3.5), CODE_BG)
example = ('## Past Editor Assignments for Similar Manuscripts\n'
           'The following are past editor assignments for\n'
           'similar manuscripts. Use these as additional context.\n'
           'Do NOT blindly copy past assignments.\n\n'
           '1. Manuscript ja-2025-00100\n'
           '   - Assigned Editor: 130958\n'
           '   - Reasoning: Expert in organic chemistry\n'
           '     with strong catalysis background...\n'
           '   - Topics: catalysis, organic chemistry\n'
           '   - Runner-up: ED002\n'
           '   - Date: 2026-02-23')
add_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.6), Inches(3.3),
             example, font_size=12, color=ACCENT_GREEN, font_name="Consolas")

# Safety note
add_rounded_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
                 RGBColor(0x2A, 0x2A, 0x44))
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
             "Safety: Memory reads are wrapped in try/except — if the store is down or search fails, the workflow continues without "
             "past assignments. Memory never breaks the main flow.",
             font_size=14, color=ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7: Files Changed
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Files Changed", font_size=32, color=WHITE, bold=True)

# New files
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "New Files  (5)", font_size=22, color=ACCENT_GREEN, bold=True)

new_files = [
    ("memory.py", "~350 lines", "Core module — checkpointer, store, save/search, format_past_assignments"),
    ("test_memory_local.py", "~430 lines", "6 unit tests — MemorySaver + InMemoryStore"),
    ("test_memory_integration.py", "~450 lines", "4 integration tests — real graph + mock LLM"),
    ("test_memory.py", "~200 lines", "4 Postgres tests (for when Docker available)"),
    ("poc_memory_demo.py", "~300 lines", "4 interactive demos for presentation"),
]

for i, (fname, size, desc) in enumerate(new_files):
    y = Inches(1.9 + i * 0.6)
    add_text_box(slide, Inches(1.0), y, Inches(2.5), Inches(0.5),
                 fname, font_size=15, color=ACCENT_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(3.5), y, Inches(1.2), Inches(0.5),
                 size, font_size=13, color=MED_GRAY)
    add_text_box(slide, Inches(4.7), y, Inches(4.5), Inches(0.5),
                 desc, font_size=13, color=LIGHT_GRAY)

# Modified files
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.5),
             "Modified Files  (4)", font_size=22, color=ORANGE, bold=True)

mod_files = [
    ("ee_graph_anthropic.py", "Accepts checkpointer/store, reads past assignments into prompt, saves to memory"),
    ("app.py", "Lifespan creates memory if POSTGRES_URI set, graceful fallback"),
    ("prompts.py", "Added {past_assignments} placeholder to V3 prompt template"),
    ("pyproject.toml", "Added langgraph-checkpoint-postgres, psycopg[binary]"),
]

for i, (fname, desc) in enumerate(mod_files):
    y = Inches(5.6 + i * 0.55)
    add_text_box(slide, Inches(1.0), y, Inches(3.0), Inches(0.5),
                 fname, font_size=15, color=ORANGE, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(4.2), y, Inches(6.5), Inches(0.5),
                 desc, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8: Test Results
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Test Results — 10/10 Tests + 4/4 Demos Passed", font_size=32, color=ACCENT_GREEN, bold=True)

# Test Suite 1
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.8),
                 RGBColor(0x2A, 0x2A, 0x44))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "Unit Tests — test_memory_local.py  (6/6  ✅)", font_size=18, color=ACCENT_GREEN, bold=True)

unit_tests = [
    "✅ test_session_memory_checkpoint_and_resume",
    "✅ test_session_memory_audit_trail",
    "✅ test_long_term_memory_store_and_retrieve",
    "✅ test_save_assignment_to_memory_helper",
    "✅ test_both_memories_combined",
    "✅ test_format_past_assignments_for_prompt",
]

for i, t in enumerate(unit_tests):
    add_text_box(slide, Inches(1.0), Inches(1.9 + i * 0.4), Inches(5.2), Inches(0.4),
                 t, font_size=14, color=LIGHT_GRAY, font_name="Consolas")

# Test Suite 2
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.8),
                 RGBColor(0x2A, 0x2A, 0x44))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(0.4),
             "Integration Tests — test_memory_integration.py  (4/4  ✅)", font_size=18, color=ACCENT_GREEN, bold=True)

int_tests = [
    "✅ test_full_workflow_with_memory",
    "    → 7+ checkpoints, assignment saved to store",
    "✅ test_resubmission_flow_with_memory",
    "    → resubmission path + 5+ checkpoints",
    "✅ test_multiple_manuscripts_build_knowledge",
    "    → 3 manuscripts, 3 records, all searchable",
    "✅ test_memory_read_injects_past_assignments",
    "    → MS-1 stored, MS-2 prompt has past data",
]

for i, t in enumerate(int_tests):
    add_text_box(slide, Inches(7.3), Inches(1.9 + i * 0.35), Inches(5.2), Inches(0.35),
                 t, font_size=13, color=LIGHT_GRAY, font_name="Consolas")

# POC Demo
add_rounded_rect(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.8),
                 RGBColor(0x2A, 0x2A, 0x44))
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
             "POC Demo — poc_memory_demo.py  (4/4  ✅)", font_size=18, color=ACCENT_GREEN, bold=True)

demos = [
    ("Demo 1: Crash Recovery", "Interrupt graph mid-run, resume from checkpoint — zero data loss"),
    ("Demo 2: Audit Trail", "Show all state changes via aget_state_history() — full visibility"),
    ("Demo 3: Store & Search", "Save 3 assignments, retrieve by key, verify data integrity"),
    ("Demo 4: Real Workflow", "Full EditorAssignmentWorkflow graph with both memories active"),
]

for i, (name, desc) in enumerate(demos):
    y = Inches(5.0 + i * 0.3)
    add_text_box(slide, Inches(1.0), y, Inches(3.0), Inches(0.3),
                 f"✅  {name}", font_size=14, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(4.5), y, Inches(7.5), Inches(0.3),
                 desc, font_size=14, color=LIGHT_GRAY)

# Note
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             "Tests use in-memory backends (same interface as Postgres). Also verified on real MSPUBS Dev RDS — 6 tables created, data ops confirmed.",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9: Why Tests Are Valid
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Why In-Memory Tests Are Valid", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "LangGraph uses a backend abstraction — like Java's JPA / Spring Data",
             font_size=18, color=LIGHT_GRAY)

# Interface diagram
add_rounded_rect(slide, Inches(1.5), Inches(1.9), Inches(4.5), Inches(0.8),
                 PURPLE, "Interface:  BaseCheckpointSaver", font_size=16, font_color=WHITE, bold=True)
add_rounded_rect(slide, Inches(7.3), Inches(1.9), Inches(4.5), Inches(0.8),
                 PURPLE, "Interface:  BaseStore", font_size=16, font_color=WHITE, bold=True)

# Test impl
add_rounded_rect(slide, Inches(0.5), Inches(3.2), Inches(3.0), Inches(1.0),
                 RGBColor(0x2A, 0x6A, 0x3A), "MemorySaver\n(in-memory)", font_size=14, font_color=WHITE)
add_text_box(slide, Inches(0.5), Inches(4.3), Inches(3.0), Inches(0.4),
             "Unit Tests  ✅", font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Prod impl
add_rounded_rect(slide, Inches(3.8), Inches(3.2), Inches(3.0), Inches(1.0),
                 RGBColor(0x1A, 0x3A, 0x5C), "AsyncPostgresSaver\n(Postgres)", font_size=14, font_color=WHITE)
add_text_box(slide, Inches(3.8), Inches(4.3), Inches(3.0), Inches(0.4),
             "Production  🚀", font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(7.3), Inches(3.2), Inches(2.8), Inches(1.0),
                 RGBColor(0x2A, 0x6A, 0x3A), "InMemoryStore\n(in-memory)", font_size=14, font_color=WHITE)
add_text_box(slide, Inches(7.3), Inches(4.3), Inches(2.8), Inches(0.4),
             "Unit Tests  ✅", font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(10.4), Inches(3.2), Inches(2.8), Inches(1.0),
                 RGBColor(0x1A, 0x3A, 0x5C), "AsyncPostgresStore\n(pgvector)", font_size=14, font_color=WHITE)
add_text_box(slide, Inches(10.4), Inches(4.3), Inches(2.8), Inches(0.4),
             "Production  🚀", font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Java analogy
add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8),
                 RGBColor(0x2A, 0x2A, 0x44))
add_text_box(slide, Inches(1.1), Inches(5.3), Inches(11), Inches(0.4),
             "Java Analogy (for the Spring Boot developers)", font_size=18, color=ORANGE, bold=True)

analogies = [
    ("JpaRepository interface", "=", "BaseCheckpointSaver / BaseStore"),
    ("H2 in-memory DB (unit tests)", "=", "MemorySaver + InMemoryStore"),
    ("Real PostgreSQL (production)", "=", "AsyncPostgresSaver + AsyncPostgresStore"),
    ("@Mock / Mockito", "=", "Monkey-patch _fetch_manuscript_data, _call_assign_api"),
]

for i, (java, eq, python) in enumerate(analogies):
    y = Inches(5.8 + i * 0.3)
    add_text_box(slide, Inches(1.3), y, Inches(4.0), Inches(0.3),
                 java, font_size=13, color=LIGHT_GRAY, font_name="Consolas")
    add_text_box(slide, Inches(5.5), y, Inches(0.5), Inches(0.3),
                 eq, font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.2), y, Inches(5.5), Inches(0.3),
                 python, font_size=13, color=ACCENT_GREEN, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10: Real RDS Verification
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Real RDS Verification — MSPUBS Dev", font_size=32, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "PostgreSQL 17.4  •  mspubs-dev.cdeku8g0y28t.us-east-1.rds.amazonaws.com  •  Schema: mspubs",
             font_size=16, color=LIGHT_GRAY)

# Tables created (left)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
             "6 Tables Created in mspubs Schema", font_size=20, color=WHITE, bold=True)

tables = [
    ("mspubs.checkpoints", "Workflow state snapshots"),
    ("mspubs.checkpoint_blobs", "Large binary state data"),
    ("mspubs.checkpoint_writes", "Write-ahead log"),
    ("mspubs.checkpoint_migrations", "Schema versioning (10 rows)"),
    ("mspubs.store", "Past assignment decisions — JSONB (3 test rows)"),
    ("mspubs.store_migrations", "Schema versioning (4 rows)"),
]

for i, (tname, desc) in enumerate(tables):
    y = Inches(2.4 + i * 0.5)
    add_text_box(slide, Inches(1.0), y, Inches(3.2), Inches(0.5),
                 f"✅  {tname}", font_size=14, color=ACCENT_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(4.3), y, Inches(3.0), Inches(0.5),
                 desc, font_size=13, color=LIGHT_GRAY)

# Operations verified (right)
add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.5),
             "Operations Verified on Real Postgres", font_size=20, color=WHITE, bold=True)

ops = [
    ("store.aput() — save assignment", "✅ PASS", ACCENT_GREEN),
    ("store.aget() — retrieve by key", "✅ PASS", ACCENT_GREEN),
    ("store.asearch() — list by namespace", "✅ PASS", ACCENT_GREEN),
    ("checkpointer.setup() — create tables", "✅ PASS", ACCENT_GREEN),
    ("pgvector extension install", "❌ Needs rds_superuser", ACCENT_RED),
]

for i, (op, status, clr) in enumerate(ops):
    y = Inches(2.4 + i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(3.8), Inches(0.5),
                 op, font_size=14, color=LIGHT_GRAY)
    add_text_box(slide, Inches(11.0), y, Inches(2.0), Inches(0.5),
                 status, font_size=14, color=clr, bold=True)

# Connection string note
add_rounded_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4), CODE_BG)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
             "Critical: Connection string must include search_path=mspubs  (no CREATE permission on public schema)",
             font_size=15, color=ORANGE, bold=True)
conn_str = 'POSTGRES_URI=postgresql://mspubs_user:<pw>@mspubs-dev...com:5432/mspubs?options=-csearch_path%3Dmspubs'
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
             conn_str, font_size=13, color=ACCENT_GREEN, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11: What's Next
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "What's Next — Remaining Work", font_size=32, color=WHITE, bold=True)

# Done column
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Completed  ✅", font_size=22, color=ACCENT_GREEN, bold=True)

done_items = [
    "memory.py — core module with Tier 2 + Tier 3",
    "Integration into ee_graph_anthropic.py + app.py",
    "Memory READ path — past assignments injected into LLM prompt",
    "6 unit tests + 4 integration tests passing",
    "4 POC demos running end-to-end",
    "Verified on real MSPUBS Dev RDS (6 tables created)",
    "Docker image built (v1.3.0-memory)",
    "Graceful fallback when no Postgres",
]

for i, item in enumerate(done_items):
    add_text_box(slide, Inches(1.0), Inches(1.8 + i * 0.45), Inches(5.5), Inches(0.45),
                 f"✅  {item}", font_size=15, color=LIGHT_GRAY)

# To-do column
add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
             "To Do  🔲", font_size=22, color=ORANGE, bold=True)

todo_items = [
    ("Install pgvector extension on RDS", "DBA action"),
    ("Push image to GHCR", "Need PAT"),
    ("Set POSTGRES_URI in K8s ConfigMap", "Config"),
    ("Deploy to dev K8s via ArgoCD", "DevOps"),
    ("End-to-end test with real manuscript", "Validation"),
    ("Retry / circuit breaker patterns", "Resilience"),
]

for i, (item, tag) in enumerate(todo_items):
    y = Inches(1.8 + i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(4.5), Inches(0.45),
                 f"🔲  {item}", font_size=15, color=LIGHT_GRAY)
    add_rounded_rect(slide, Inches(11.5), y + Inches(0.05), Inches(1.3), Inches(0.35),
                     RGBColor(0x3A, 0x3A, 0x5A), tag, font_size=10, font_color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11: Summary / Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Key Takeaways", font_size=36, color=WHITE, bold=True)

takeaways = [
    ("1", "Crash Recovery", "If a pod restarts mid-workflow, no work is lost — resumes from last checkpoint", ACCENT_BLUE),
    ("2", "Full Audit Trail", "Every node execution is recorded — debug any decision, anytime", ACCENT_BLUE),
    ("3", "Learning from History", "Completed assignments are stored with vector embeddings — semantic search over past decisions", ACCENT_GREEN),
    ("4", "Zero Breaking Changes", "Memory is opt-in via POSTGRES_URI — without it, everything works exactly as before", ORANGE),
    ("5", "Production Ready Pattern", "Same LangGraph interface for in-memory (tests) and Postgres (prod) — validated by 10 tests + 4 demos + real RDS", PURPLE),
]

for i, (num, title, desc, color) in enumerate(takeaways):
    y = Inches(1.6 + i * 1.05)
    add_rounded_rect(slide, Inches(0.8), y, Inches(0.7), Inches(0.7),
                     color, num, font_size=24, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y, Inches(10), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.4), Inches(10), Inches(0.5),
                 desc, font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════════════════

output_path = "docs/memory-implementation-poc-v2.pptx"
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
